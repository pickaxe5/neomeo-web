"""docs/frontend-to-backend-requests.md #14: 업로드된 기획 문서 파일에서 텍스트를 추출한다.
원본 바이너리는 저장하지 않고 추출한 텍스트만 project_documents.content에 남긴다 — "원본 파일
저장 여부" 확인 필요 항목은 이 방향(텍스트만 보관)으로 결정."""

from io import BytesIO

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

MAX_UPLOAD_SIZE_BYTES = 10 * 1024 * 1024  # 10MB
SUPPORTED_EXTENSIONS = (".pdf", ".docx", ".pptx", ".txt", ".md")


class DocumentExtractionError(Exception):
    pass


def _extract_pdf(data: bytes) -> str:
    reader = PdfReader(BytesIO(data))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(data: bytes) -> str:
    doc = DocxDocument(BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_pptx(data: bytes) -> str:
    prs = Presentation(BytesIO(data))
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs)
                if text:
                    lines.append(text)
    return "\n".join(lines)


def _extract_plain_text(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".txt": _extract_plain_text,
    ".md": _extract_plain_text,
}


def extract_text(filename: str, data: bytes) -> str:
    """filename 확장자에 맞는 추출기로 텍스트를 뽑는다.

    지원하지 않는 확장자이거나, 추출 결과가 비어있으면(스캔본 이미지 PDF 등 텍스트 레이어가
    없는 파일) DocumentExtractionError를 올린다 — "파싱 실패 시 에러 응답" 확인 필요 항목에
    대한 답: 400 + 사람이 읽을 수 있는 메시지."""
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        raise DocumentExtractionError(
            f"지원하지 않는 파일 형식입니다 ({ext or '확장자 없음'}). "
            f"지원 형식: {', '.join(SUPPORTED_EXTENSIONS)}"
        )

    try:
        text = extractor(data)
    except Exception as exc:  # noqa: BLE001 — 파싱 라이브러리별 예외를 하나로 묶어 일관되게 응답
        raise DocumentExtractionError(f"파일을 읽는 중 오류가 발생했습니다: {exc}") from exc

    text = text.strip()
    if not text:
        raise DocumentExtractionError(
            "파일에서 텍스트를 추출하지 못했습니다 (스캔본 이미지 PDF처럼 텍스트가 없는 "
            "파일일 수 있습니다)."
        )
    return text
